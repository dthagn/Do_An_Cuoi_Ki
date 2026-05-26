"""One-off generator for reports/report.docx and reports/slides.pptx."""
from __future__ import annotations

from pathlib import Path

BASE = Path(__file__).resolve().parent.parent
REPORTS = BASE / "reports"
IMAGES = BASE / "images"

METRICS = [
    ("Logistic Regression", "0.8436", "0.7879", "0.8725"),
    ("Decision Tree", "0.8268", "0.7634", "0.8502"),
    ("Random Forest", "0.8101", "0.7424", "0.8449"),
]
BEST_MODEL = "Logistic Regression"
RF_PARAMS = "max_depth=4, n_estimators=100"


def build_docx(path: Path) -> None:
    from docx import Document
    from docx.shared import Inches, Pt

    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    doc.add_heading("Đồ án: Dự đoán sống sót Titanic", 0)
    doc.add_paragraph(
        "Báo cáo tóm tắt pipeline Machine Learning — từ khám phá dữ liệu đến "
        "mô hình triển khai và giải thích (SHAP)."
    )

    doc.add_heading("1. Giới thiệu", level=1)
    doc.add_paragraph(
        "Bài toán phân loại nhị phân trên tập Kaggle Titanic: dự đoán hành khách "
        "có sống sót (Survived = 1) hay không. Dữ liệu huấn luyện gồm 891 mẫu; "
        "tập test Kaggle 418 mẫu (không nhãn)."
    )

    doc.add_heading("2. Khám phá & tiền xử lý", level=1)
    doc.add_paragraph(
        "Tỷ lệ sống sót trên tập train khoảng 38,4% — lớp mất cân bằng nhẹ. "
        "Các bước làm sạch: điền Age theo median theo Title; Embarked theo mode; "
        "Fare theo median. Feature engineering: Family_Size, Is_Alone, Title (từ Name), "
        "Deck (từ Cabin), Group_Size (từ Ticket)."
    )
    path_survival = IMAGES / "survival_rate.png"
    if path_survival.exists():
        doc.add_picture(str(path_survival), width=Inches(5.0))
        doc.add_paragraph("Hình 1 — Phân bố nhãn Survived.", style="Caption")

    doc.add_heading("3. Mô hình & đánh giá", level=1)
    doc.add_paragraph(
        "Ba mô hình so sánh trên hold-out 20% (stratified): Logistic Regression "
        "(chuẩn hóa số), Decision Tree và Random Forest (GridSearchCV, F1). "
        f"Mô hình nộp Kaggle và lưu model.pkl: {BEST_MODEL} (F1 cao nhất trên hold-out)."
    )
    table = doc.add_table(rows=1, cols=4)
    hdr = table.rows[0].cells
    hdr[0].text = "Mô hình"
    hdr[1].text = "Accuracy"
    hdr[2].text = "F1"
    hdr[3].text = "ROC-AUC"
    for name, acc, f1, roc in METRICS:
        row = table.add_row().cells
        row[0].text = name
        row[1].text = acc
        row[2].text = f1
        row[3].text = roc

    for img_name, caption in [
        ("roc_curve.png", "Hình 2 — Đường cong ROC (các mô hình)."),
        ("confusion_matrix.png", f"Hình 3 — Ma trận nhầm lẫn ({BEST_MODEL})."),
        ("feature_importance.png", "Hình 4 — Feature importance (Random Forest)."),
    ]:
        img = IMAGES / img_name
        if img.exists():
            doc.add_picture(str(img), width=Inches(5.2))
            doc.add_paragraph(caption, style="Caption")

    doc.add_heading("4. Giải thích mô hình (SHAP)", level=1)
    doc.add_paragraph(
        f"Random Forest sau tinh chỉnh ({RF_PARAMS}) dùng cho phân tích SHAP toàn cục "
        "(TreeExplainer trên tập hold-out). SHAP cho thấy Sex, Pclass, Fare, Title "
        "thường đóng vai trò quan trọng — phù hợp với câu chuyện lịch sử Titanic "
        "(ưu tiên phụ nữ/trẻ em, hạng vé)."
    )
    shap_img = IMAGES / "shap_summary.png"
    if shap_img.exists():
        doc.add_picture(str(shap_img), width=Inches(5.5))
        doc.add_paragraph("Hình 5 — SHAP summary plot (Random Forest).", style="Caption")

    doc.add_heading("5. Triển khai", level=1)
    doc.add_paragraph(
        "submission.csv: dự đoán 418 hành khách test. model.pkl: pipeline sklearn "
        "(preprocess + classifier) cho app Streamlit. Chạy: pip install -r requirements.txt; "
        "streamlit run app.py."
    )

    doc.add_heading("6. Kết luận & hạn chế", level=1)
    doc.add_paragraph(
        "Đã xây dựng pipeline end-to-end, tránh leakage bằng sklearn Pipeline. "
        "Hạn chế: dữ liệu nhỏ; một số feature (Ticket, Cabin) thưa; chưa calibration "
        "xác suất. Hướng mở rộng: cross-validation lồng nhau, ensemble, giải thích SHAP "
        "theo từng dự đoán (local) trên app."
    )

    doc.save(path)


def build_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def title_slide(title: str, subtitle: str = "") -> None:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

    def bullet_slide(title: str, bullets: list[str]) -> None:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, line in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = Pt(20)

    def image_slide(title: str, img_path: Path, note: str = "") -> None:
        layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if img_path.exists():
            slide.shapes.add_picture(
                str(img_path), Inches(0.8), Inches(1.3), width=Inches(11.5)
            )
        if note:
            tx = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5))
            tx.text_frame.text = note

    title_slide(
        "Dự đoán sống sót Titanic",
        "Machine Learning — Kaggle Titanic\nĐồ án cuối kỳ",
    )

    bullet_slide(
        "Vấn đề",
        [
            "Dự đoán Survived (0/1) cho hành khách Titanic",
            "891 mẫu train · 418 mẫu test (Kaggle)",
            "Mục tiêu: pipeline ML + giải thích + demo",
        ],
    )

    bullet_slide(
        "Dữ liệu & insight EDA",
        [
            "~38% sống sót trên train — mất cân bằng nhẹ",
            "Thiếu Age, Cabin; categorical Sex, Embarked, Pclass",
            "Engineering: Title, Deck, Family_Size, Group_Size",
        ],
    )

    image_slide(
        "EDA — Survival rate",
        IMAGES / "survival_rate.png",
        "Phân bố nhãn trên tập huấn luyện",
    )

    bullet_slide(
        "Phương pháp",
        [
            "sklearn Pipeline: preprocess + model (tránh leakage)",
            "Hold-out 20% stratified · metric: F1, ROC-AUC",
            "GridSearchCV (Decision Tree, Random Forest)",
        ],
    )

    bullet_slide(
        "Kết quả so sánh (hold-out)",
        [
            f"Logistic Regression — F1 {METRICS[0][2]} (chọn deploy)",
            f"Decision Tree — F1 {METRICS[1][2]}",
            f"Random Forest — F1 {METRICS[2][2]} (dùng SHAP)",
        ],
    )

    image_slide("ROC curves", IMAGES / "roc_curve.png", "So sánh khả năng phân tách lớp")
    image_slide(
        "Confusion matrix",
        IMAGES / "confusion_matrix.png",
        f"Mô hình triển khai: {BEST_MODEL}",
    )
    image_slide(
        "Feature importance",
        IMAGES / "feature_importance.png",
        "Random Forest — hướng giải thích bổ sung",
    )
    image_slide(
        "SHAP — Random Forest",
        IMAGES / "shap_summary.png",
        "Giải thích toàn cục: feature nào đẩy dự đoán sống/chết",
    )

    bullet_slide(
        "Demo Streamlit",
        [
            "Nhập thông tin hành khách → xác suất sống sót",
            "Hiển thị SHAP summary (train sẵn từ notebook)",
            "model.pkl + images/shap_summary.png",
        ],
    )

    bullet_slide(
        "Kết luận",
        [
            "Pipeline hoàn chỉnh: EDA → train → Kaggle submission",
            f"Deploy: {BEST_MODEL}; giải thích: Random Forest + SHAP",
            "Hạn chế: dữ liệu nhỏ — hướng ensemble / local SHAP",
        ],
    )

    title_slide("Cảm ơn", "Q&A")

    prs.save(path)


def main() -> None:
    REPORTS.mkdir(exist_ok=True)
    docx_path = REPORTS / "report.docx"
    pptx_path = REPORTS / "slides.pptx"
    build_docx(docx_path)
    build_pptx(pptx_path)
    print(f"Wrote: {docx_path}")
    print(f"Wrote: {pptx_path}")


if __name__ == "__main__":
    main()
